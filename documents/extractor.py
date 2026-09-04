from pathlib import Path

import pymupdf
from docx import Document

from app.config import settings
from openpyxl import load_workbook
from pptx import Presentation


def extract_text(file_path: Path) -> str:
    extension = file_path.suffix.lower()

    if extension == ".pdf":
        return _extract_pdf(file_path)

    if extension == ".docx":
        return _extract_docx(file_path)

    if extension == ".xlsx":
        return _extract_xlsx(file_path)

    if extension == ".pptx":
        return _extract_pptx(file_path)

    if extension == ".txt":
        return file_path.read_text(encoding="utf-8", errors="ignore")

    raise ValueError(f"Unsupported file type: {extension}")


def _extract_pdf(file_path: Path) -> str:
    pages = []

    with pymupdf.open(file_path) as document:
        for page in document:
            pages.append(page.get_text())

    return "\n".join(pages).strip()


def extract_pdf_visuals(file_path: Path, file_id: str, max_pages: int = 20) -> list[dict[str, str | int]]:
    """Render only PDF pages containing images, with deterministic local artifact IDs."""
    if file_path.suffix.lower() != ".pdf" or max_pages <= 0:
        return []
    output_dir = settings.data_dir / "visuals" / file_id
    artifacts: list[dict[str, str | int]] = []
    try:
        with pymupdf.open(file_path) as document:
            for page_number, page in enumerate(document, start=1):
                if len(artifacts) >= max_pages:
                    break
                if not page.get_images(full=True):
                    continue
                artifact_id = f"visuals/{file_id}/{file_id}-p{page_number:04d}.png"
                target = output_dir / f"{file_id}-p{page_number:04d}.png"
                target.parent.mkdir(parents=True, exist_ok=True)
                page.get_pixmap(matrix=pymupdf.Matrix(1, 1), alpha=False).save(target)
                artifacts.append({"image_id": artifact_id, "page": page_number})
    except (RuntimeError, OSError, ValueError):
        return []
    return artifacts


def _extract_docx(file_path: Path) -> str:
    document = Document(file_path)

    paragraphs = [
        paragraph.text
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]

    return "\n".join(paragraphs).strip()


def _extract_xlsx(file_path: Path) -> str:
    workbook = load_workbook(
        filename=file_path,
        read_only=True,
        data_only=True,
    )

    lines = []

    for worksheet in workbook.worksheets:
        lines.append(f"[Sheet: {worksheet.title}]")

        for row in worksheet.iter_rows(values_only=True):
            values = [
                str(value)
                for value in row
                if value is not None
            ]

            if values:
                lines.append(" | ".join(values))

    workbook.close()

    return "\n".join(lines).strip()


def _extract_pptx(file_path: Path) -> str:
    presentation = Presentation(file_path)

    slides = []

    for index, slide in enumerate(presentation.slides, start=1):
        slides.append(f"[Slide {index}]")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slides.append(shape.text.strip())

    return "\n".join(slides).strip()


